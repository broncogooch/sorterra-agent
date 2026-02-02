#!/usr/bin/env python3
"""
ACME Corp Full-Spectrum File Generator
Generates 23 types of business-themed files for Sorterra testing.
"""

import os
import random
import string
import json
import xml.etree.ElementTree as ET
import sqlite3
import wave
import struct
import zipfile
import csv
import io
import time
from pathlib import Path
from datetime import datetime, timedelta
from dotenv import load_dotenv

# Third-party imports
from PIL import Image, ImageDraw, ImageFont
from reportlab.lib.pagesizes import letter, A4
from reportlab.pdfgen import canvas
from reportlab.lib.colors import HexColor
from openpyxl import Workbook
from openpyxl.styles import Font, Fill, PatternFill
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from docx import Document
from docx.shared import Inches as DocxInches, Pt as DocxPt
import pandas as pd

# =============================================================================
# BUSINESS CONTEXT DATA
# =============================================================================

BUSINESS_NAME = "ACME Corp"
DEPARTMENTS = ["Finance", "Engineering", "Marketing", "HR", "Operations", "Legal", "Sales", "IT", "Product"]
PROJECTS = ["Alpha", "Beta", "Gamma", "Delta", "Omega", "Phoenix", "Starlight", "Nova"]
DOC_TYPES = ["Invoice", "Project Plan", "Meeting Notes", "Contract", "Marketing Strategy", "Financial Report", "Budget", "Technical Spec", "User Manual", "Audit Report"]
VENDORS = ["AWS", "Google Cloud", "Microsoft", "Stripe", "GitHub", "Slack", "Zoom", "Twilio"]
EMPLOYEES = ["Alice Smith", "Bob Jones", "Charlie Brown", "Diana Prince", "Edward Norton", "Fiona Gallagher"]
TOPICS = ["Cloud Migration", "Database Upgrade", "Q3 Earnings", "Security Audit", "Client Presentation"]

# =============================================================================
# THEMED UTILITIES
# =============================================================================

def random_sentence():
    return f"The {random.choice(DEPARTMENTS)} team is finalized on {random.choice(PROJECTS)} for {random.choice(TOPICS)}."

def random_paragraph(n=3):
    return " ".join([random_sentence() for _ in range(n)])

def random_filename():
    styles = [
        lambda: f"Project_{random.choice(PROJECTS)}_{random.choice(DOC_TYPES).replace(' ', '_')}",
        lambda: f"{random.choice(VENDORS)}_Invoice_{random.randint(1000, 9999)}",
        lambda: f"{random.choice(DEPARTMENTS)}_Update_{datetime.now().strftime('%Y%m%d')}",
    ]
    return random.choice(styles)()

def random_color():
    return (random.randint(0, 255), random.randint(0, 255), random.randint(0, 255))

# =============================================================================
# 23 BUSINESS-THEMED GENERATORS
# =============================================================================

def generate_txt(fp, size=1):
    with open(fp, 'w') as f:
        f.write(f"{BUSINESS_NAME} INTERNAL MEMO\n{random_paragraph(10)}")

def generate_csv(fp, size=1):
    with open(fp, 'w', newline='') as f:
        w = csv.writer(f)
        w.writerow(["ID", "Project", "Department", "Cost Center"])
        for i in range(20):
            w.writerow([f"ACME-{i}", random.choice(PROJECTS), random.choice(DEPARTMENTS), random.randint(100, 900)])

def generate_json(fp, size=1):
    data = {"company": BUSINESS_NAME, "active_projects": PROJECTS, "audit_log": [random_sentence() for _ in range(5)]}
    with open(fp, 'w') as f:
        json.dump(data, f, indent=2)

def generate_xml(fp, size=1):
    root = ET.Element("BusinessData", company=BUSINESS_NAME)
    for p in PROJECTS:
        ET.SubElement(root, "Project", name=p, status="Active")
    ET.ElementTree(root).write(fp)

def generate_html(fp, size=1):
    content = f"<html><body><h1>{BUSINESS_NAME} Dashboard</h1><p>{random_paragraph()}</p></body></html>"
    with open(fp, 'w') as f: f.write(content)

def generate_md(fp, size=1):
    with open(fp, 'w') as f:
        f.write(f"# Project {random.choice(PROJECTS)} Roadmap\n\n## Overview\n{random_paragraph()}")

def generate_png(fp, size=1):
    img = Image.new('RGB', (800, 400), random_color())
    draw = ImageDraw.Draw(img)
    draw.text((100, 150), f"{BUSINESS_NAME}: {random.choice(PROJECTS)}", fill=(255,255,255))
    img.save(fp, 'PNG')

def generate_jpg(fp, size=1):
    img = Image.new('RGB', (800, 400), random_color())
    img.save(fp, 'JPEG')

def generate_gif(fp, size=1):
    frames = [Image.new('RGB', (400, 400), random_color()) for _ in range(5)]
    frames[0].save(fp, save_all=True, append_images=frames[1:], duration=200, loop=0)

def generate_bmp(fp, size=1):
    Image.new('RGB', (400, 400), random_color()).save(fp, 'BMP')

def generate_svg(fp, size=1):
    with open(fp, 'w') as f:
        f.write(f'<svg xmlns="http://www.w3.org/2000/svg" width="200" height="200"><rect width="100%" height="100%" fill="blue"/><text x="10" y="50" fill="white">{BUSINESS_NAME}</text></svg>')

def generate_pdf(fp, size=1):
    c = canvas.Canvas(fp, pagesize=letter)
    c.drawString(100, 750, f"{BUSINESS_NAME} - CONFIDENTIAL")
    c.drawString(100, 730, f"Subject: Project {random.choice(PROJECTS)} Analysis")
    c.save()

def generate_xlsx(fp, size=1):
    wb = Workbook()
    ws = wb.active
    ws.append(["Dept", "Project", "Manager"])
    ws.append([random.choice(DEPARTMENTS), random.choice(PROJECTS), random.choice(EMPLOYEES)])
    wb.save(fp)

def generate_pptx(fp, size=1):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = f"{BUSINESS_NAME} Strategic Review"
    prs.save(fp)

def generate_docx(fp, size=1):
    doc = Document()
    doc.add_heading(f"Project {random.choice(PROJECTS)} Master Plan", 0)
    doc.add_paragraph(random_paragraph())
    doc.save(fp)

def generate_sqlite(fp, size=1):
    conn = sqlite3.connect(fp)
    conn.execute("CREATE TABLE Projects (id INT, name TEXT, dept TEXT)")
    conn.execute(f"INSERT INTO Projects VALUES (1, '{random.choice(PROJECTS)}', '{random.choice(DEPARTMENTS)}')")
    conn.commit()
    conn.close()

def generate_parquet(fp, size=1):
    df = pd.DataFrame({"Project": PROJECTS, "Code": [f"ACME-{i}" for i in range(len(PROJECTS))]})
    df.to_parquet(fp)

def generate_wav(fp, size=1):
    with wave.open(fp, 'w') as wf:
        wf.setnchannels(1); wf.setsampwidth(2); wf.setframerate(44100)
        for _ in range(1000): wf.writeframes(struct.pack('<h', random.randint(-32000, 32000)))

def generate_zip(fp, size=1):
    with zipfile.ZipFile(fp, 'w') as zf:
        zf.writestr("manifest.txt", f"Archive for {BUSINESS_NAME} - Project {random.choice(PROJECTS)}")

def generate_log(fp, size=1):
    with open(fp, 'w') as f:
        for _ in range(20): f.write(f"[{datetime.now()}] INFO: {random.choice(DEPARTMENTS)} access granted.\n")

def generate_yaml(fp, size=1):
    with open(fp, 'w') as f:
        f.write(f"company: {BUSINESS_NAME}\ndepartment: {random.choice(DEPARTMENTS)}\nactive: true")

def generate_ini(fp, size=1):
    with open(fp, 'w') as f:
        f.write(f"[System]\nCompany={BUSINESS_NAME}\nProject={random.choice(PROJECTS)}")

def generate_rtf(fp, size=1):
    with open(fp, 'w') as f:
        f.write(r"{\rtf1\ansi " + f"{BUSINESS_NAME} - {random.choice(DOC_TYPES)}" + r"}")

# =============================================================================
# MAIN LOGIC
# =============================================================================

FILE_GENERATORS = {
    'txt': generate_txt, 'csv': generate_csv, 'json': generate_json, 'xml': generate_xml,
    'html': generate_html, 'md': generate_md, 'png': generate_png, 'jpg': generate_jpg,
    'gif': generate_gif, 'bmp': generate_bmp, 'svg': generate_svg, 'pdf': generate_pdf,
    'xlsx': generate_xlsx, 'pptx': generate_pptx, 'docx': generate_docx, 'sqlite': generate_sqlite,
    'parquet': generate_parquet, 'wav': generate_wav, 'zip': generate_zip, 'log': generate_log,
    'yaml': generate_yaml, 'ini': generate_ini, 'rtf': generate_rtf
}

def main():
    load_dotenv()
    # Ensure this points to Sorterra's data folder
    output_path = Path("./data/test_folder")
    output_path.mkdir(parents=True, exist_ok=True)
    
    num_files = 23 # Generates one of each for perfect testing
    for i, (ext, func) in enumerate(FILE_GENERATORS.items()):
        fname = f"{random_filename()}.{ext}"
        func(str(output_path / fname))
        print(f"[{i+1}/{num_files}] Created themed file: {fname}")

if __name__ == "__main__":
    main()